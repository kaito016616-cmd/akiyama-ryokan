# -*- coding: utf-8 -*-
"""
卒研中間報告書1 発表スライド作成スクリプト
テーマ: ABSがもたらす影響かつ簡易ABSの作成
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- カラー・フォント設定（明るめ・ゴシック体） ----------
FONT_NAME = "游ゴシック"
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_MAIN = RGBColor(0x1E, 0x88, 0xE5)      # 明るい青（アクセント）
COLOR_SUB = RGBColor(0x43, 0xA0, 0x47)       # 明るい緑（グラウンドイメージ）
COLOR_DARK = RGBColor(0x26, 0x2A, 0x2E)      # 本文用ダークグレー
COLOR_LIGHT_BG = RGBColor(0xF2, 0xF8, 0xFC)  # 薄い水色背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY = RGBColor(0x8A, 0x8F, 0x94)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def set_font(run, size, color=COLOR_DARK, bold=False, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name
    # 日本語フォントを明示的に指定
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def add_bg(slide, color=COLOR_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_textbox(slide, left, top, width, height, text, size, color=COLOR_DARK,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name=FONT_NAME,
                 line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_font(run, size, color, bold, name)
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=COLOR_DARK,
                 bold_first=False, name=FONT_NAME, space_after=10, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = f"● {item}"
        set_font(run, size, color, False, name)
    return box


def add_footer(slide, page_num, total, title="ABSがもたらす影響かつ簡易ABSの作成"):
    add_textbox(slide, Inches(0.4), Inches(7.12), Inches(8), Inches(0.35),
                title, 10, COLOR_GRAY)
    add_textbox(slide, Inches(12.3), Inches(7.12), Inches(0.7), Inches(0.35),
                f"{page_num}/{total}", 10, COLOR_GRAY, align=PP_ALIGN.RIGHT)


def add_header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_MAIN
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, 28, COLOR_WHITE, True)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12), Inches(0.4),
                    subtitle, 14, COLOR_SUB, bold=True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), Inches(13.333), Pt(4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_SUB
    accent.line.fill.background()
    accent.shadow.inherit = False


def add_photo_placeholder(slide, left, top, width, height, label="写真挿入予定"):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT_BG
    box.line.color.rgb = COLOR_MAIN
    box.line.width = Pt(1.5)
    box.line.dash_style = None
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {label} ]"
    set_font(run, 14, COLOR_MAIN, True)
    return box


def add_box_arrow_flow(slide, top, items, box_w=Inches(1.75), box_h=Inches(0.95),
                        gap=Inches(0.25), start_left=Inches(0.5), size=12):
    left = start_left
    for i, item in enumerate(items):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_MAIN if i % 2 == 0 else COLOR_SUB
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item
        set_font(run, size, COLOR_WHITE, True)
        left_next = left + box_w
        if i < len(items) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, left_next, top + box_h / 2 - Inches(0.12),
                gap, Inches(0.24))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_GRAY
            arrow.line.fill.background()
            arrow.shadow.inherit = False
        left = left_next + gap


TOTAL_SLIDES = 12
page = 0


def next_page():
    global page
    page += 1
    return page


# ============================================================
# 1. 表紙
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
top_band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(4.6))
top_band.fill.solid()
top_band.fill.fore_color.rgb = COLOR_LIGHT_BG
top_band.line.fill.background()
top_band.shadow.inherit = False

accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.6), SLIDE_W, Pt(6))
accent.fill.solid()
accent.fill.fore_color.rgb = COLOR_SUB
accent.line.fill.background()
accent.shadow.inherit = False

add_textbox(s, Inches(1), Inches(1.3), Inches(11.3), Inches(0.5),
            "令和8年度 卒業研究 中間報告書1", 20, COLOR_MAIN, bold=True)
add_textbox(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.8),
            "ABSがもたらす影響かつ\n簡易ABSの作成", 40, COLOR_DARK, bold=True, line_spacing=1.2)

add_photo_placeholder(s, Inches(9.6), Inches(1.3), Inches(3.0), Inches(3.0), "野球/ABS写真")

add_textbox(s, Inches(1), Inches(5.1), Inches(6), Inches(0.4),
            "情報メディア工学科／データサイエンス学科", 15, COLOR_DARK)
add_textbox(s, Inches(1), Inches(5.55), Inches(10), Inches(0.9),
            "提出者：稲木陽秋　神田蒼生　鶴身敬太　黛祐貴　宮嶋海音", 15, COLOR_DARK)
add_textbox(s, Inches(1), Inches(6.05), Inches(6), Inches(0.4),
            "指導教員：北久保茂 准教授", 15, COLOR_DARK)
add_textbox(s, Inches(1), Inches(6.5), Inches(6), Inches(0.4),
            "提出日：2026年7月17日", 13, COLOR_GRAY)

# ============================================================
# 2. 目次
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "目次")
items = [
    "研究背景・目的",
    "本研究の有用性",
    "ABSとは／疑似ABSの仕組み",
    "課題",
    "システム設計",
    "研究の進捗（文献調査・技術調査・データ収集）",
    "比較・考察",
    "今後の日程計画",
]
left = Inches(0.9)
top = Inches(1.6)
for i, item in enumerate(items):
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
    num.fill.solid()
    num.fill.fore_color.rgb = COLOR_MAIN if i % 2 == 0 else COLOR_SUB
    num.line.fill.background()
    num.shadow.inherit = False
    tf = num.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(i + 1)
    set_font(run, 16, COLOR_WHITE, True)
    add_textbox(s, left + Inches(0.7), top + Inches(0.03), Inches(9.5), Inches(0.5),
                item, 17, COLOR_DARK)
    top += Inches(0.66)
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 3. 研究背景・目的
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "研究背景・目的")
add_textbox(s, Inches(0.6), Inches(1.5), Inches(3), Inches(0.4), "背景", 18, COLOR_MAIN, bold=True)
add_bullets(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(3.5), [
    "プロ野球でABS（自動ボール・ストライク判定システム）が導入",
    "判定の公平性が向上した一方、高性能機器が必要で一般利用は困難",
    "画像処理技術を用いた「疑似ABS」を開発し、日常の野球遊びで\n手軽に利用できるシステムを目指す",
], size=17)

add_textbox(s, Inches(7.0), Inches(1.5), Inches(3), Inches(0.4), "目的", 18, COLOR_SUB, bold=True)
add_bullets(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(3.5), [
    "ABSを参考に、画像処理を用いた疑似ABSを開発する",
    "キャッチボールやストラックアウトなど日常の遊びで\n判定を楽しめるシステムを実現する",
    "精度・実用性が確認できれば草野球への応用も視野",
], size=17)

add_photo_placeholder(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.4), "ABS導入シーン等の写真")
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 4. 本研究の有用性
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "本研究の有用性")
add_bullets(s, Inches(0.7), Inches(1.7), Inches(11.8), Inches(2.5), [
    "高価な設備を必要とせず、スマートフォンなど身近な機器でボール判定が可能に",
    "日常の野球遊びをより楽しめるだけでなく、判定の補助ツールとしても活用できる",
    "精度が向上すれば、草野球など審判の負担軽減・判定の公平性向上にも貢献",
], size=19, space_after=16)

# 効果を示す簡単な図解（3つのメリットカード）
cards = [
    ("低コスト", "スマホだけで判定OK"),
    ("手軽さ", "日常の遊びで活用可能"),
    ("将来性", "草野球への応用も視野"),
]
left = Inches(0.9)
w = Inches(3.7)
gap = Inches(0.4)
for i, (title, desc) in enumerate(cards):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.5), w, Inches(1.9))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = COLOR_SUB
    card.line.width = Pt(1.5)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    set_font(run, 20, COLOR_MAIN, True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = desc
    set_font(run2, 14, COLOR_DARK)
    left += w + gap
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 5. ABSとは／疑似ABSの仕組み
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "ABSとは／疑似ABSの仕組み")

add_textbox(s, Inches(0.6), Inches(1.45), Inches(5.9), Inches(0.4), "ABS（本家）", 18, COLOR_MAIN, bold=True)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(5.9), Inches(2.3), [
    "複数台の高速度カメラでボール軌道を追跡",
    "三次元空間上で軌道を推定し判定",
    "判定の公平性・正確性を向上",
], size=15)

add_textbox(s, Inches(6.9), Inches(1.45), Inches(5.9), Inches(0.4), "本研究の疑似ABS", 18, COLOR_SUB, bold=True)
add_bullets(s, Inches(6.9), Inches(1.9), Inches(5.9), Inches(2.3), [
    "スマートフォンで撮影した動画を利用",
    "画像処理でボール位置を検出・追跡",
    "軌道がストライクゾーンを通過したか判定",
], size=15)

# 対比図解
box1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.5), Inches(4.8), Inches(1.6))
box1.fill.solid(); box1.fill.fore_color.rgb = COLOR_MAIN; box1.line.fill.background(); box1.shadow.inherit = False
tf = box1.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "高速度カメラ×複数台\n(高コスト・高精度)"
set_font(run, 16, COLOR_WHITE, True)

vs = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.05), Inches(4.95), Inches(0.7), Inches(0.7))
vs.fill.solid(); vs.fill.fore_color.rgb = COLOR_WHITE; vs.line.color.rgb = COLOR_GRAY; vs.shadow.inherit = False
tf = vs.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "VS"
set_font(run, 16, COLOR_GRAY, True)

box2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(4.5), Inches(4.8), Inches(1.6))
box2.fill.solid(); box2.fill.fore_color.rgb = COLOR_SUB; box2.line.fill.background(); box2.shadow.inherit = False
tf = box2.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "スマートフォン1台\n(低コスト・手軽)"
set_font(run, 16, COLOR_WHITE, True)

add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 6. 課題
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "課題")
add_bullets(s, Inches(0.7), Inches(1.7), Inches(11.8), Inches(3.5), [
    "スマートフォンのカメラは高速度カメラより性能が低く、ボール高速移動時の検出が困難",
    "照明条件・背景・撮影位置・撮影角度の違いにより判定精度が変化する",
    "ストライクゾーンの設定方法や、ボールの軌道推定精度についても検討が必要",
], size=19, space_after=18)
add_photo_placeholder(s, Inches(0.7), Inches(4.6), Inches(11.8), Inches(2.0), "検出困難な場面の写真例")
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 7. システム設計（全体像・6ステップ）
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "システム設計（全体構成）")
add_textbox(s, Inches(0.6), Inches(1.4), Inches(11), Inches(0.4),
            "以下の6つの処理を順番に実行し、ストライク・ボールを判定する", 16, COLOR_DARK)

steps1 = ["動画入力", "キャリブレーション", "ボール検出"]
steps2 = ["軌跡フィルタリング\n・通過位置推定", "ストライク判定", "結果表示"]
add_box_arrow_flow(s, Inches(2.1), steps1, box_w=Inches(3.6), box_h=Inches(1.1),
                    gap=Inches(0.35), start_left=Inches(1.0), size=15)
add_box_arrow_flow(s, Inches(3.6), steps2, box_w=Inches(3.6), box_h=Inches(1.1),
                    gap=Inches(0.35), start_left=Inches(1.0), size=13)

add_bullets(s, Inches(0.7), Inches(5.1), Inches(11.8), Inches(1.9), [
    "動画入力：側面・後方の2方向から撮影した映像を使用",
    "ボール検出：背景差分法(MOG2)＋輪郭の円形度評価で検出",
    "ストライク判定：側面カメラでY軸、後方カメラでX軸を独立評価し、両軸ゾーン内でストライク",
], size=14, space_after=6)
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 8. 研究の進捗①（類似研究・画像処理技術）
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "研究の進捗①　文献・技術調査")

add_textbox(s, Inches(0.6), Inches(1.5), Inches(5.9), Inches(0.4), "類似研究の調査", 18, COLOR_MAIN, bold=True)
add_bullets(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(2.3), [
    "複数カメラでボール位置を取得し、三次元軌道を\n推定する手法が一般的",
    "画像処理と深層学習の組み合わせで高精度な\n判定が実現可能と把握",
], size=15)

add_textbox(s, Inches(6.9), Inches(1.5), Inches(5.9), Inches(0.4), "画像処理技術の調査", 18, COLOR_SUB, bold=True)
add_bullets(s, Inches(6.9), Inches(2.0), Inches(5.9), Inches(2.3), [
    "OpenCVによる輪郭検出・色抽出を調査",
    "YOLOなどの物体検出技術も調査",
    "今後、本研究に適した手法を選定し実装予定",
], size=15)

add_textbox(s, Inches(0.6), Inches(4.6), Inches(11.8), Inches(0.4), "判定システムの調査", 18, COLOR_MAIN, bold=True)
add_bullets(s, Inches(0.6), Inches(5.05), Inches(11.8), Inches(1.5), [
    "ストライクゾーンを基準に、ボール軌道が通過したかで判定する仕組みを確認",
    "スマートフォン映像でも利用可能な簡易判定アルゴリズムを検討中",
], size=15)
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 9. 研究の進捗②（動画データの収集）
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "研究の進捗②　動画データの収集")
add_bullets(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(3.5), [
    "YouTube上の野球動画・判定場面を収集",
    "自分たちで実際にピッチングを行い、\nストライク・ボール判定が確認できる\n場面を撮影・切り出し",
    "今後のプログラム開発・判定精度評価に\n利用できるデータセットを作成中",
], size=17, space_after=14)
add_photo_placeholder(s, Inches(7.0), Inches(1.7), Inches(5.6), Inches(4.7), "ピッチング撮影の様子")
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 10. 比較・考察
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "比較・考察")

# 表：既存ABS vs 疑似ABS
rows, cols = 3, 3
left, top, width, height = Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.4)
table_shape = s.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table
table.columns[0].width = Inches(3.3)
table.columns[1].width = Inches(4.3)
table.columns[2].width = Inches(4.3)

headers = ["", "既存ABS", "本研究（疑似ABS）"]
data = [
    ["設備コスト", "高い（複数台の高速度カメラ）", "低い（スマートフォン1台）"],
    ["判定精度", "非常に高い", "既存より低下する可能性あり"],
]
for c, text in enumerate(headers):
    cell = table.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_MAIN
    cell.text_frame.word_wrap = True
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, 15, COLOR_WHITE, True)

for r, row in enumerate(data, start=1):
    for c, text in enumerate(row):
        cell = table.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_LIGHT_BG if r % 2 == 1 else COLOR_WHITE
        cell.text_frame.word_wrap = True
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        set_font(run, 14, COLOR_DARK, bold=(c == 0))

add_bullets(s, Inches(0.7), Inches(4.4), Inches(11.9), Inches(2.4), [
    "導入コストを大幅に抑えられる点が本研究の特徴",
    "一方で判定精度は既存システムより低下する可能性があるため、画像処理手法や\n判定アルゴリズムの改善が必要",
    "今後は異なる撮影条件・ボール速度で実験を行い、判定精度や動作の安定性を評価予定",
], size=15, space_after=10)
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 11. 今後の日程計画
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "今後の日程計画")

rows, cols = 6, 2
left, top, width, height = Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.3)
table_shape = s.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(8.4)

sched = [
    ("時期", "作業内容"),
    ("〜8月", "検証実験、動画データ収集の継続"),
    ("9月〜10月", "データ集計、システム評価・改善"),
    ("10月〜11月", "プログラム改良、評価の継続"),
    ("12月〜1月", "卒業論文作成"),
    ("1月〜2月", "確認、発表資料作成"),
]
for r, (a, b) in enumerate(sched):
    header_row = (r == 0)
    for c, text in enumerate([a, b]):
        cell = table.cell(r, c)
        cell.fill.solid()
        if header_row:
            cell.fill.fore_color.rgb = COLOR_MAIN
        else:
            cell.fill.fore_color.rgb = COLOR_LIGHT_BG if r % 2 == 1 else COLOR_WHITE
        cell.text_frame.word_wrap = True
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        set_font(run, 15, COLOR_WHITE if header_row else COLOR_DARK, bold=header_row or c == 0)
add_footer(s, next_page(), TOTAL_SLIDES)

# ============================================================
# 12. まとめ
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header_bar(s, "まとめ")
add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.8), Inches(3.5), [
    "スマートフォンのみで利用できる疑似ABSの開発を進めている",
    "文献調査・技術調査を通じてシステム設計の方向性を決定した",
    "今後はプログラム実装と検証実験を進め、判定精度・実用性を評価していく",
], size=20, space_after=18)

thanks = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.3), Inches(5.4), Inches(4.7), Inches(1.2))
thanks.fill.solid()
thanks.fill.fore_color.rgb = COLOR_SUB
thanks.line.fill.background()
thanks.shadow.inherit = False
tf = thanks.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "ご清聴ありがとうございました"
set_font(run, 20, COLOR_WHITE, True)
add_footer(s, next_page(), TOTAL_SLIDES)

# ---------- 保存 ----------
out_path = r"C:\Users\kaito\Documents\claudecodetest\presentation\ABS中間報告スライド.pptx"
prs.save(out_path)
print("saved:", out_path)
